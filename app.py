import streamlit as st
import streamlit_shadcn_ui as ui
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
import io
import os
import time
import math
import json

# GitHub 라이브러리
try:
    from github import Github
    GITHUB_AVAILABLE = True
except ImportError:
    GITHUB_AVAILABLE = False

# --- 설정 ---
TEMPLATE_FILE = "template.pptx"
SIDEBAR_LOGO = "assets/bossgolf.svg"
LOGO_DIR = "assets/logos"
ARTWORK_DIR = "assets/artworks"
CSS_FILE = "style.css"
ARTWORK_META_FILE = os.path.join(ARTWORK_DIR, "_meta.json")

# --- 텍스트 좌표/스타일 스펙 (mm 기준) ---
TEXT_SPECS = {
    "season": {
        "left": 22.5,
        "top": 12.5,
        "width": 83.33,
        "height": 9.49,
        "font_name": "Averta PE Extrabold",
        "font_size": 12,
        "bold": True,
        "color_hex": "#000000",  # 사용자 선택으로 덮어씀
    },
    "category": {
        "left": 9.5,
        "top": 24.1,
        "width": 117.05,
        "height": 13.85,
        "font_name": "Averta PE Extrabold",
        "font_size": 24,
        "bold": True,
        "color_hex": "#987147",
    },
    "code": {
        "left": 9.5,
        "top": 32.5,
        "width": 117.05,
        "height": 13.85,
        "font_name": "Averta PE Extrabold",
        "font_size": 24,
        "bold": True,
        "color_hex": "#000000",
    },
    # page number is managed by slide master placeholder
}

COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM = 169.9
COLORWAY_TWO_ITEMS_LABEL_TOP_MM = 114.8
COLORWAY_TWO_ITEMS_LABEL_GAP_MM = 28.0
COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM = 169.9
COLORWAY_THREE_ITEMS_LABEL_TOP_MM = 114.8
COLORWAY_THREE_ITEMS_LABEL_GAP_MM = 28.0
COLORWAY_IMAGE_WIDTH_MM = 27.0
COLORWAY_IMAGE_TOP_MM = 120.0
MAIN_IMAGE_CENTER_X_MM = 65.0
MAIN_IMAGE_CENTER_Y_MM = 94.3
MAIN_IMAGE_WIDTH_MM = 90.0
LOGO_CENTER_X_MM = 148.4
LOGO_CENTER_Y_MM = 53.9
LOGO_HEIGHT_MM = 23.7
ARTWORK_CENTER_X_MM = 148.4
ARTWORK_START_TOP_MM = 77.2
ARTWORK_DEFAULT_WIDTH_MM = 30.0
ARTWORK_PORTRAIT_HEIGHT_MM = 20.0
ARTWORK_SMALL_WIDTH_MM = 12.0
ARTWORK_VERTICAL_GAP_MM = 5.0
ARTWORK_MODE_DEFAULT = "default"
ARTWORK_MODE_HORIZONTAL = "horizontal"
ARTWORK_MODE_SMALL = "small"
ARTWORK_MODE_LABEL_TO_VALUE = {
    "기본": ARTWORK_MODE_DEFAULT,
    "가로 타입": ARTWORK_MODE_HORIZONTAL,
    "작은 아트워크": ARTWORK_MODE_SMALL,
}
ARTWORK_MODE_VALUE_TO_LABEL = {v: k for k, v in ARTWORK_MODE_LABEL_TO_VALUE.items()}

# --- 유틸리티 함수 ---
def init_folders():
    for folder in [LOGO_DIR, ARTWORK_DIR]:
        if not os.path.exists(folder): os.makedirs(folder)

def load_artwork_meta():
    if not os.path.exists(ARTWORK_META_FILE):
        return {}
    try:
        with open(ARTWORK_META_FILE, "r", encoding="utf-8") as f:
            data = json.load(f)
            return data if isinstance(data, dict) else {}
    except Exception:
        return {}

def save_artwork_meta(meta):
    with open(ARTWORK_META_FILE, "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)

def get_artwork_mode(art_name, meta=None):
    if meta is None:
        meta = load_artwork_meta()
    mode = meta.get(art_name, ARTWORK_MODE_DEFAULT)
    if mode not in (ARTWORK_MODE_DEFAULT, ARTWORK_MODE_HORIZONTAL, ARTWORK_MODE_SMALL):
        return ARTWORK_MODE_DEFAULT
    return mode

def load_css(file_name):
    if os.path.exists(file_name):
        with open(file_name) as f:
            st.markdown(f'<style>{f.read()}</style>', unsafe_allow_html=True)

def get_files(folder_path):
    if not os.path.exists(folder_path): return []
    return [f for f in os.listdir(folder_path) if f.lower().endswith(('.png', '.jpg', '.jpeg', '.svg'))]

def get_layout_by_matching_name(prs, targets):
    target_set = {t.strip().lower() for t in targets}
    for layout in prs.slide_layouts:
        matching_name = (layout._element.get("matchingName") or "").strip().lower()
        if matching_name in target_set:
            return layout
    return None

def get_layout_by_name(prs, targets):
    target_set = {t.strip().lower() for t in targets}
    for layout in prs.slide_layouts:
        if (layout.name or "").strip().lower() in target_set:
            return layout
    return None

def get_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").strip()

def find_layout_anchor(layout):
    anchors = {}
    text_shapes = []
    empty_shapes = []

    for shp in layout.shapes:
        txt = get_text(shp)
        if txt:
            text_shapes.append(shp)
        elif getattr(shp, "shape_type", None) is not None:
            empty_shapes.append(shp)

    for shp in text_shapes:
        txt = get_text(shp).upper()
        if "LOGO" in txt:
            anchors["logo_label"] = shp
        if "ARTWORK" in txt:
            anchors["artwork_label"] = shp
        if "RRP" in txt:
            anchors["rrp_label"] = shp
        if "COLORWAY" in txt:
            anchors["color_label"] = shp

    # 텍스트 라벨 근처의 빈 도형을 이미지 박스로 사용
    def nearest_empty(label_key):
        label = anchors.get(label_key)
        if not label:
            return None
        lx = label.left + (label.width // 2)
        ly = label.top + (label.height // 2)
        candidates = []
        for shp in empty_shapes:
            w, h = shp.width, shp.height
            if w <= 0 or h <= 0:
                continue
            # 너무 작은 가이드/점은 제외
            if w < Mm(20) or h < Mm(10):
                continue
            sx = shp.left + (w // 2)
            sy = shp.top + (h // 2)
            dist = abs(sx - lx) + abs(sy - ly)
            candidates.append((dist, shp))
        return min(candidates, key=lambda x: x[0])[1] if candidates else None

    anchors["logo_box"] = nearest_empty("logo_label")
    anchors["artwork_box"] = nearest_empty("artwork_label")
    return anchors

def find_text_slot(layout, keywords):
    keyset = [k.upper() for k in keywords]
    for shp in layout.shapes:
        txt = get_text(shp).upper()
        if not txt:
            continue
        if any(k in txt for k in keyset):
            return shp
    return None

def hex_to_rgbcolor(hex_color):
    if not hex_color:
        return None
    v = hex_color.strip().lstrip("#")
    if len(v) != 6:
        return None
    try:
        return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except:
        return None

def replace_text_in_slide(slide, tokens, value, color_override=None):
    token_set = [t.upper() for t in tokens]
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        current = (shp.text or "")
        up = current.upper()
        if any(t in up for t in token_set):
            style = extract_text_style(shp)
            shp.text_frame.clear()
            shp.text_frame.text = value
            p = shp.text_frame.paragraphs[0]
            apply_text_style(p, style)
            if color_override is not None:
                try:
                    p.font.color.rgb = color_override
                except:
                    pass
            return True
    return False

def extract_text_style(shape):
    style = {"align": None, "size": None, "bold": None, "italic": None, "name": None, "color_rgb": None}
    if not getattr(shape, "has_text_frame", False):
        return style
    tf = shape.text_frame
    if not tf.paragraphs:
        return style
    p = tf.paragraphs[0]
    style["align"] = p.alignment
    f = p.font
    style["size"] = f.size
    style["bold"] = f.bold
    style["italic"] = f.italic
    style["name"] = f.name
    try:
        if f.color and getattr(f.color, "rgb", None):
            style["color_rgb"] = f.color.rgb
    except:
        pass
    return style

def apply_text_style(paragraph, style):
    if style.get("align") is not None:
        paragraph.alignment = style["align"]
    f = paragraph.font
    if style.get("size") is not None:
        f.size = style["size"]
    if style.get("bold") is not None:
        f.bold = style["bold"]
    if style.get("italic") is not None:
        f.italic = style["italic"]
    if style.get("name"):
        f.name = style["name"]
    if style.get("color_rgb") is not None:
        try:
            f.color.rgb = style["color_rgb"]
        except:
            pass

def add_styled_textbox(slide, slot_shape, text, fallback_left, fallback_top, fallback_w, fallback_h, fallback_size):
    if slot_shape is not None:
        tb = slide.shapes.add_textbox(slot_shape.left, slot_shape.top, slot_shape.width, slot_shape.height)
        style = extract_text_style(slot_shape)
    else:
        tb = slide.shapes.add_textbox(fallback_left, fallback_top, fallback_w, fallback_h)
        style = {"size": fallback_size, "bold": True}
    tb.text_frame.text = text
    p = tb.text_frame.paragraphs[0]
    apply_text_style(p, style)
    return tb

def add_text_by_spec(slide, text, spec, color_override=None):
    tb = slide.shapes.add_textbox(
        Mm(spec["left"]),
        Mm(spec["top"]),
        Mm(spec["width"]),
        Mm(spec["height"]),
    )
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = spec["font_name"]
    f.size = Pt(spec["font_size"])
    f.bold = spec.get("bold")

    # 일부 PowerPoint 환경에서 font.name만으로는 폰트 유지가 불안정해서
    # run 레벨의 latin/ea/cs 타입페이스를 함께 강제 지정한다.
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rPr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rPr.append(node)
        node.set("typeface", spec["font_name"])

    color = color_override or spec.get("color_hex")
    rgb = hex_to_rgbcolor(color)
    if rgb is not None:
        try:
            f.color.rgb = rgb
        except:
            pass
    return tb

def add_text_at(slide, text, left_mm, top_mm, width_mm, height_mm, font_name, font_size_pt, color_hex="#000000", bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(font_size_pt)
    f.bold = bold
    rgb = hex_to_rgbcolor(color_hex)
    if rgb is not None:
        f.color.rgb = rgb
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rPr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rPr.append(node)
        node.set("typeface", font_name)
    return tb

def format_color_name(name):
    if not name:
        return ""
    return str(name).strip().upper()

def has_slide_number_placeholder(slide):
    for shp in slide.shapes:
        if not getattr(shp, "is_placeholder", False):
            continue
        try:
            if shp.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER:
                return True
        except:
            continue
    return False

def strip_vendor_watermark(prs):
    markers = ("VORLAGENBAUER", "ERSTELLT DURCH")
    targets = list(prs.slide_masters)
    for m in prs.slide_masters:
        targets.extend(list(m.slide_layouts))

    for container in targets:
        for shp in container.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            txt = (shp.text or "").upper()
            if any(marker in txt for marker in markers):
                shp.text_frame.clear()

def ensure_slide_number_enabled(prs):
    """
    Force the presentation-level header/footer flag so slide-number placeholders
    on layouts are shown without requiring 'Apply to All' in PowerPoint UI.
    """
    hf = prs._element.find(qn("p:hf"))
    if hf is None:
        hf = OxmlElement("p:hf")
        prs._element.append(hf)
    hf.set("sldNum", "1")

# --- 깃허브 연동 ---
def get_github_repo():
    if not GITHUB_AVAILABLE: return None
    try:
        return Github(st.secrets["github"]["token"]).get_repo(st.secrets["github"]["repo_name"])
    except: return None

def upload_file(file_obj, folder_path):
    with open(os.path.join(folder_path, file_obj.name), "wb") as f:
        f.write(file_obj.getbuffer())
    repo = get_github_repo()
    if repo:
        try:
            path = f"{folder_path}/{file_obj.name}"
            content = file_obj.getvalue()
            branch = st.secrets["github"].get("branch", "main")
            try:
                contents = repo.get_contents(path, ref=branch)
                repo.update_file(path, f"Update {file_obj.name}", content, contents.sha, branch=branch)
            except:
                repo.create_file(path, f"Upload {file_obj.name}", content, branch=branch)
            return True
        except: return False
    return True

def delete_file_asset(filename, folder_path):
    local = os.path.join(folder_path, filename)
    if os.path.exists(local): os.remove(local)
    repo = get_github_repo()
    if repo:
        try:
            path = f"{folder_path}/{filename}"
            branch = st.secrets["github"].get("branch", "main")
            contents = repo.get_contents(path, ref=branch)
            repo.delete_file(path, f"Delete {filename}", contents.sha, branch=branch)
        except: pass

# --- PPT 생성 로직 ---
def create_pptx(products):
    if os.path.exists(TEMPLATE_FILE): prs = Presentation(TEMPLATE_FILE)
    else: prs = Presentation()
    strip_vendor_watermark(prs)
    ensure_slide_number_enabled(prs)

    # 레이아웃 선택 우선순위:
    # 1) matchingName = default (요청사항 우선)
    # 2) matchingName = title
    # 3) 표시 이름 = HB Title / Content, CUSTOM
    # 3) 기존 인덱스 폴백
    selected_layout = (
        get_layout_by_matching_name(prs, ["default"])
        or get_layout_by_matching_name(prs, ["title"])
        or get_layout_by_name(prs, ["HB Title / Content", "CUSTOM"])
    )
    if selected_layout is None:
        selected_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    layout_anchors = find_layout_anchor(selected_layout)
    artwork_meta = load_artwork_meta()

    for data in products:
        slide = prs.slides.add_slide(selected_layout)
        season_color_hex = data.get("season_color") or TEXT_SPECS["season"]["color_hex"]

        # 텍스트: 좌표 고정 매핑
        season_name = data.get("season_item", "")
        if season_name:
            add_text_by_spec(
                slide,
                season_name,
                TEXT_SPECS["season"],
                color_override=season_color_hex,
            )

        add_text_by_spec(
            slide,
            data["name"],
            TEXT_SPECS["category"],
        )
        add_text_by_spec(
            slide,
            data["code"],
            TEXT_SPECS["code"],
        )

        # RRP (표시만 함)
        if data.get('rrp'):
            rrp_left = layout_anchors["rrp_label"].left if layout_anchors.get("rrp_label") else Mm(250)
            rrp_top = layout_anchors["rrp_label"].top if layout_anchors.get("rrp_label") else Mm(15)
            rrp = slide.shapes.add_textbox(rrp_left, rrp_top, Mm(50), Mm(15))
            rrp.text_frame.text = f"RRP : {data['rrp']}"
            rrp.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # 메인 이미지
        if data['main_image']:
            main_pic = slide.shapes.add_picture(
                data['main_image'],
                left=Mm(0),
                top=Mm(0),
                width=Mm(MAIN_IMAGE_WIDTH_MM),
            )
            main_pic.left = int(Mm(MAIN_IMAGE_CENTER_X_MM) - (main_pic.width / 2))
            main_pic.top = int(Mm(MAIN_IMAGE_CENTER_Y_MM) - (main_pic.height / 2))
        
        # 로고
        if data['logo'] and data['logo'] != "선택 없음":
            p_logo = os.path.join(LOGO_DIR, data['logo'])
            if os.path.exists(p_logo):
                logo_pic = slide.shapes.add_picture(
                    p_logo,
                    left=Mm(0),
                    top=Mm(0),
                    height=Mm(LOGO_HEIGHT_MM),
                )
                logo_pic.left = int(Mm(LOGO_CENTER_X_MM) - (logo_pic.width / 2))
                logo_pic.top = int(Mm(LOGO_CENTER_Y_MM) - (logo_pic.height / 2))
        
        # 아트워크 (여러 개 선택 가능, 위에서 아래로 스택)
        if data['artworks']:
            current_top = int(Mm(ARTWORK_START_TOP_MM))
            gap_emu = int(Mm(ARTWORK_VERTICAL_GAP_MM))
            for art_item in data['artworks']:
                art_name = art_item.get("name") if isinstance(art_item, dict) else art_item
                p_art = os.path.join(ARTWORK_DIR, art_name)
                if not os.path.exists(p_art):
                    continue
                artwork_mode = get_artwork_mode(art_name, artwork_meta)
                if artwork_mode == ARTWORK_MODE_SMALL:
                    art_pic = slide.shapes.add_picture(
                        p_art,
                        left=Mm(0),
                        top=Mm(0),
                        width=Mm(ARTWORK_SMALL_WIDTH_MM),
                    )
                elif artwork_mode == ARTWORK_MODE_HORIZONTAL:
                    art_pic = slide.shapes.add_picture(
                        p_art,
                        left=Mm(0),
                        top=Mm(0),
                        width=Mm(ARTWORK_DEFAULT_WIDTH_MM),
                    )
                else:
                    art_pic = slide.shapes.add_picture(
                        p_art,
                        left=Mm(0),
                        top=Mm(0),
                        height=Mm(ARTWORK_PORTRAIT_HEIGHT_MM),
                    )
                art_pic.left = int(Mm(ARTWORK_CENTER_X_MM) - (art_pic.width / 2))
                art_pic.top = current_top
                current_top += art_pic.height + gap_emu

        # 컬러웨이
        sx, sy, w, g = 180, COLORWAY_IMAGE_TOP_MM, COLORWAY_IMAGE_WIDTH_MM, 5
        color_count = len(data.get('colors', []))
        is_two_item_single_row = color_count == 2
        is_three_item_single_row = color_count == 3
        if layout_anchors.get("color_label"):
            # 라벨 아래 우하단 영역을 시작점으로 사용
            sx = layout_anchors["color_label"].left / 36000.0
            sy = COLORWAY_IMAGE_TOP_MM
        # 2개/1줄은 라벨/이미지 기준점을 완전히 고정한다.
        if is_two_item_single_row:
            sx = COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM
            sy = COLORWAY_IMAGE_TOP_MM
        # 3개/1줄도 고정 좌표 규칙을 사용한다.
        if is_three_item_single_row:
            sx = COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM
            sy = COLORWAY_IMAGE_TOP_MM
        per_row = 3
        row_gap = 8
        img_h = 30
        rows = max(1, math.ceil(len(data['colors']) / per_row)) if data.get('colors') else 1
        circled_nums = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]
        for i, c in enumerate(data['colors']):
            row = i // per_row
            col = i % per_row
            # 2줄 이상이면 아래줄 고정 후 위로 쌓기
            cy = sy - (rows - 1 - row) * (img_h + row_gap + 10)
            cx = sx + (col * (w + g))
            if is_two_item_single_row:
                cx = COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_TWO_ITEMS_LABEL_GAP_MM)
                cy = COLORWAY_IMAGE_TOP_MM
            if is_three_item_single_row:
                cx = COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_THREE_ITEMS_LABEL_GAP_MM)
                cy = COLORWAY_IMAGE_TOP_MM
            if c['img']:
                slide.shapes.add_picture(c['img'], left=Mm(cx), top=Mm(cy), width=Mm(COLORWAY_IMAGE_WIDTH_MM))
            # 1줄/2개 케이스: 지정 좌표에서 ①CAMEL 형식으로 라벨 표기
            if is_two_item_single_row and rows == 1:
                label = f"{circled_nums[i]}{format_color_name(c.get('name'))}"
                add_text_at(
                    slide=slide,
                    text=label,
                    left_mm=COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_TWO_ITEMS_LABEL_GAP_MM),
                    top_mm=COLORWAY_TWO_ITEMS_LABEL_TOP_MM,
                    width_mm=32.0,
                    height_mm=5.0,
                    font_name="Averta Light",
                    font_size_pt=10,
                    color_hex="#000000",
                    bold=False,
                    align=PP_ALIGN.LEFT,
                )
            elif is_three_item_single_row and rows == 1:
                label = f"{circled_nums[i]}{format_color_name(c.get('name'))}"
                add_text_at(
                    slide=slide,
                    text=label,
                    left_mm=COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_THREE_ITEMS_LABEL_GAP_MM),
                    top_mm=COLORWAY_THREE_ITEMS_LABEL_TOP_MM,
                    width_mm=32.0,
                    height_mm=5.0,
                    font_name="Averta Light",
                    font_size_pt=10,
                    color_hex="#000000",
                    bold=False,
                    align=PP_ALIGN.LEFT,
                )
            else:
                tb = slide.shapes.add_textbox(Mm(cx), Mm(cy + img_h + 2), Mm(w), Mm(10))
                tb.text_frame.text = format_color_name(c.get('name'))
                tb.text_frame.paragraphs[0].font.size = Pt(9)
                tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# =========================================================
# APP MAIN
# =========================================================
st.set_page_config(page_title="BOSS Golf Admin", layout="wide", initial_sidebar_state="expanded")
init_folders()
load_css(CSS_FILE)

if 'product_list' not in st.session_state:
    st.session_state.product_list = []
if 'colorway_items' not in st.session_state:
    st.session_state.colorway_items = []
if 'colorway_sig' not in st.session_state:
    st.session_state.colorway_sig = []
if 'artwork_items' not in st.session_state:
    st.session_state.artwork_items = []

# --- 1. 좌측 사이드바 ---
with st.sidebar:
    if os.path.exists(SIDEBAR_LOGO):
        st.image(SIDEBAR_LOGO, width=140)
    else:
        st.header("BOSS Golf")
    
    st.markdown("---")
    selected_menu = ui.tabs(
        options=['슬라이드 제작', '로고&아트워크 관리'],
        defaultValue='슬라이드 제작',
        key="sidebar_menu",
    )

    # [수정] 깃허브 연동 상태 표시 제거됨


# --- 2. 메인 콘텐츠 ---

if selected_menu == '슬라이드 제작':
    st.title("슬라이드 제작")
    st.markdown("제품 정보를 입력하여 스펙 시트를 생성합니다.")
    st.markdown("---")

    tab_editor, tab_queue = st.tabs(["정보 입력", "생성 대기열"])
    
    # 탭 1: 입력
    with tab_editor:
        st.subheader("1. 기본 정보")
        c1, c2 = st.columns([3, 1])
        with c1:
            season_item = st.text_input("시즌 아이템명", "JETSET LUXE")
            season_color = st.color_picker("시즌 텍스트 색상", "#000000")
            p_name = st.text_input("제품명", "MEN'S T-SHIRTS")
            p_code = st.text_input("품번 (필수)", placeholder="예: BKFTM1581")
        
        st.markdown("<br>", unsafe_allow_html=True)
        st.subheader("2. 디자인 자산")
        
        main_img = st.file_uploader("메인 이미지", type=['png','jpg'], help="슬라이드 좌측에 크게 들어갈 이미지")
        if main_img:
            st.caption("메인 이미지 미리보기")
            st.image(main_img, width=220)
        
        c3, c4 = st.columns(2)
        with c3:
            logo_options = ["선택 없음"] + sorted(get_files(LOGO_DIR), key=str.casefold)
            s_logo = st.selectbox("로고 선택", logo_options)
            if s_logo != "선택 없음":
                st.caption("로고 미리보기")
                st.image(os.path.join(LOGO_DIR, s_logo), width=180)
        
        with c4:
            available_artworks = sorted(get_files(ARTWORK_DIR), key=str.casefold)
            selected_names = st.multiselect(
                "아트워크 선택",
                available_artworks,
                default=[item["name"] for item in st.session_state.artwork_items if item["name"] in available_artworks],
                key="artwork_multiselect",
            )

            selected_set = set(selected_names)
            current_names = [item["name"] for item in st.session_state.artwork_items]
            current_set = set(current_names)
            if current_set != selected_set:
                kept = [item for item in st.session_state.artwork_items if item["name"] in selected_set]
                kept_names = {item["name"] for item in kept}
                added_names = sorted([n for n in selected_names if n not in kept_names], key=str.casefold)
                added = [{"name": n} for n in added_names]
                st.session_state.artwork_items = kept + added

            if st.session_state.artwork_items:
                st.caption("아트워크 미리보기 및 순서")
                for idx, item in enumerate(st.session_state.artwork_items):
                    row = st.columns([1.2, 4, 0.8, 0.8])
                    with row[0]:
                        st.image(os.path.join(ARTWORK_DIR, item["name"]), width=70)
                    with row[1]:
                        st.write(item["name"])
                    with row[2]:
                        if st.button("↑", key=f"art_up_{item['name']}", disabled=(idx == 0)):
                            st.session_state.artwork_items[idx - 1], st.session_state.artwork_items[idx] = (
                                st.session_state.artwork_items[idx],
                                st.session_state.artwork_items[idx - 1],
                            )
                            st.rerun()
                    with row[3]:
                        if st.button("↓", key=f"art_down_{item['name']}", disabled=(idx == len(st.session_state.artwork_items) - 1)):
                            st.session_state.artwork_items[idx + 1], st.session_state.artwork_items[idx] = (
                                st.session_state.artwork_items[idx],
                                st.session_state.artwork_items[idx + 1],
                            )
                            st.rerun()
            else:
                st.caption("선택된 아트워크 없음")

            selected_artworks = [item["name"] for item in st.session_state.artwork_items]

        st.markdown("---")
        st.subheader("3. 컬러웨이 (Colorways)")

        uploaded_colors = st.file_uploader(
            "컬러웨이 이미지 업로드 (최대 4개)",
            type=['png','jpg'],
            accept_multiple_files=True,
            key="colorway_uploader"
        )

        current_sig = []
        if uploaded_colors:
            current_sig = [(f.name, getattr(f, "size", None)) for f in uploaded_colors[:4]]

        if current_sig != st.session_state.colorway_sig:
            st.session_state.colorway_items = []
            for i, f in enumerate(uploaded_colors[:4] if uploaded_colors else []):
                file_id = f"{f.name}:{getattr(f, 'size', 0)}:{i}"
                st.session_state.colorway_items.append({"id": file_id, "file": f, "name": ""})
            st.session_state.colorway_sig = current_sig

        colors_input = []
        if st.session_state.colorway_items:
            st.info("업로드가 완료되었습니다. 색상명을 입력하고 필요하면 순서를 변경하세요.")
            for idx, item in enumerate(st.session_state.colorway_items):
                row = st.columns([1, 4, 1, 1])
                with row[0]:
                    st.image(item["file"], width=72)
                with row[1]:
                    name_val = st.text_input(
                        f"색상명 {idx+1}",
                        value=item["name"],
                        key=f"cw_name_{item['id']}"
                    )
                    st.session_state.colorway_items[idx]["name"] = name_val
                with row[2]:
                    if st.button("↑", key=f"cw_up_{item['id']}", disabled=(idx == 0)):
                        st.session_state.colorway_items[idx - 1], st.session_state.colorway_items[idx] = (
                            st.session_state.colorway_items[idx],
                            st.session_state.colorway_items[idx - 1],
                        )
                        st.rerun()
                with row[3]:
                    if st.button("↓", key=f"cw_down_{item['id']}", disabled=(idx == len(st.session_state.colorway_items) - 1)):
                        st.session_state.colorway_items[idx + 1], st.session_state.colorway_items[idx] = (
                            st.session_state.colorway_items[idx],
                            st.session_state.colorway_items[idx + 1],
                        )
                        st.rerun()

            colors_input = [{"img": item["file"], "name": item["name"]} for item in st.session_state.colorway_items]

        st.markdown("<br>", unsafe_allow_html=True)

        if st.button("대기열에 추가", type="primary"):
            if not p_code or not main_img:
                st.error("품번과 메인 이미지는 필수입니다.")
            else:
                st.session_state.product_list.append({
                    "season_item": season_item,
                    "season_color": season_color,
                    "name":p_name, "code":p_code, "rrp":"", 
                    "main_image":main_img, 
                    "logo":s_logo, 
                    "artworks": selected_artworks,
                    "artwork": selected_artworks[0] if selected_artworks else "선택 없음",
                    "colors": colors_input
                })
                st.toast(f"'{p_code}' 대기열에 추가됨")
                st.success(f"'{p_code}' 추가 완료!")

    # 탭 2: 대기열
    with tab_queue:
        c_head, c_btn = st.columns([4, 1])
        with c_head: st.subheader(f"생성 대기 목록 ({len(st.session_state.product_list)})")
        with c_btn:
            if st.button("목록 비우기"):
                st.session_state.product_list = []
                st.rerun()
        
        if not st.session_state.product_list:
            st.info("대기 중인 항목이 없습니다.")
        else:
            for idx, item in enumerate(st.session_state.product_list):
                with st.expander(f"{idx+1}. {item['code']} - {item['name']}"):
                    cols = st.columns([1, 4])
                    cols[0].image(item['main_image'])
                    artwork_names = []
                    for art in item.get('artworks', []):
                        if isinstance(art, dict):
                            artwork_names.append(art.get("name", ""))
                        else:
                            artwork_names.append(str(art))
                    art_str = ", ".join([n for n in artwork_names if n]) if artwork_names else "-"
                    cols[1].write(f"컬러: {len(item['colors'])}개 | 로고: {item['logo']} | 아트워크: {art_str}")
            
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("PPT 생성 및 다운로드", type="primary"):
                try:
                    ppt = create_pptx(st.session_state.product_list)
                    st.download_button("PPT 다운로드 (.pptx)", ppt, "BOSS_Golf_SpecSheet.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
                except Exception as e:
                    st.error(f"PPT 생성 실패: {e}")


elif selected_menu == '로고&아트워크 관리':
    st.title("자산 관리")
    st.markdown("디자인 자산을 업로드하고 관리합니다.")
    st.markdown("---")
    
    active_tab = ui.tabs(options=['로고', '아트워크'], defaultValue='로고', key="asset_tabs")
    target_dir = LOGO_DIR if active_tab == '로고' else ARTWORK_DIR
    if active_tab == '아트워크':
        st.info(
            "아트워크 타입 안내\n"
            "- 기본: 높이 20mm (너비 자동)\n"
            "- 가로 타입: 너비 30mm (높이 자동) - 가로로 긴 형태의 아트워크인 경우 선택해주세요\n"
            "- 작은 아트워크: 너비 12mm (높이 자동) - 사이즈가 작은 아트워크인 경우 선택"
        )
    
    st.subheader(f"{active_tab} 업로드")
    uploaded = st.file_uploader("파일 업로드", type=['png','jpg','svg'], accept_multiple_files=True)
    if uploaded and st.button("저장하기"):
        with st.spinner("저장 중..."):
            for f in uploaded: upload_file(f, target_dir)
            if active_tab == '아트워크':
                artwork_meta = load_artwork_meta()
                for f in uploaded:
                    artwork_meta.setdefault(f.name, ARTWORK_MODE_DEFAULT)
                save_artwork_meta(artwork_meta)
        st.success("완료")
        time.sleep(1)
        st.rerun()
    
    st.markdown("---")
    st.subheader("보유 파일 목록")
    files = get_files(target_dir)
    artwork_meta = load_artwork_meta()
    
    if not files:
        st.info("파일이 없습니다.")
    else:
        if active_tab == '아트워크':
            files = sorted(files, key=str.casefold)
            changed = False
            for f in files:
                if f not in artwork_meta:
                    artwork_meta[f] = ARTWORK_MODE_DEFAULT
                    changed = True
            stale_keys = [k for k in artwork_meta.keys() if k not in files]
            for k in stale_keys:
                artwork_meta.pop(k, None)
                changed = True
            if changed:
                save_artwork_meta(artwork_meta)

        cols = st.columns(5)
        for i, f in enumerate(files):
            with cols[i%5]:
                st.image(os.path.join(target_dir, f), use_container_width=True)
                st.caption(f)
                if active_tab == '아트워크':
                    current_mode = get_artwork_mode(f, artwork_meta)
                    current_label = ARTWORK_MODE_VALUE_TO_LABEL.get(current_mode, "기본")
                    selected_label = st.radio(
                        "타입",
                        options=["기본", "가로 타입", "작은 아트워크"],
                        index=["기본", "가로 타입", "작은 아트워크"].index(current_label),
                        key=f"art_mode_{f}",
                    )
                    selected_mode = ARTWORK_MODE_LABEL_TO_VALUE[selected_label]
                    if selected_mode != current_mode:
                        artwork_meta[f] = selected_mode
                        save_artwork_meta(artwork_meta)
                if st.button("삭제", key=f"del_{f}"):
                    delete_file_asset(f, target_dir)
                    if active_tab == '아트워크':
                        artwork_meta.pop(f, None)
                        save_artwork_meta(artwork_meta)
                    time.sleep(1)
                    st.rerun()
