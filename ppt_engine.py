from __future__ import annotations

import io
import json
import os
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Mm, Pt

# Text specs (mm)
TEXT_SPECS = {
    "season": {
        "left": 22.5,
        "top": 12.5,
        "width": 83.33,
        "height": 9.49,
        "font_name": "Averta PE Extrabold",
        "font_size": 12,
        "bold": True,
        "color_hex": "#000000",
    },
    "category": {
        "left": 9.5,
        "top": 24.1,
        "width": 117.05,
        "height": 13.85,
        "font_name": "Averta PE Extrabold",
        "font_size": 24,
        "bold": True,
        "color_hex": "#987147",
    },
    "code": {
        "left": 9.5,
        "top": 32.5,
        "width": 117.05,
        "height": 13.85,
        "font_name": "Averta PE Extrabold",
        "font_size": 24,
        "bold": True,
        "color_hex": "#000000",
    },
}

COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM = 169.9
COLORWAY_TWO_ITEMS_LABEL_TOP_MM = 114.8
COLORWAY_TWO_ITEMS_LABEL_GAP_MM = 28.0
COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM = 169.9
COLORWAY_THREE_ITEMS_LABEL_TOP_MM = 114.8
COLORWAY_THREE_ITEMS_LABEL_GAP_MM = 28.0
COLORWAY_IMAGE_WIDTH_MM = 27.0
COLORWAY_IMAGE_TOP_MM = 120.0

MAIN_IMAGE_CENTER_X_MM = 65.0
MAIN_IMAGE_CENTER_Y_MM = 94.3
MAIN_IMAGE_WIDTH_MM = 90.0

LOGO_CENTER_X_MM = 148.4
LOGO_CENTER_Y_MM = 53.9
LOGO_HEIGHT_MM = 23.7

ARTWORK_CENTER_X_MM = 148.4
ARTWORK_START_TOP_MM = 77.2
ARTWORK_DEFAULT_WIDTH_MM = 30.0
ARTWORK_PORTRAIT_HEIGHT_MM = 20.0
ARTWORK_SMALL_WIDTH_MM = 12.0
ARTWORK_VERTICAL_GAP_MM = 5.0

ARTWORK_MODE_DEFAULT = "default"
ARTWORK_MODE_HORIZONTAL = "horizontal"
ARTWORK_MODE_SMALL = "small"


def _hex_to_rgbcolor(hex_color: str | None):
    if not hex_color:
        return None
    v = hex_color.strip().lstrip("#")
    if len(v) != 6:
        return None
    try:
        return RGBColor(int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    except Exception:
        return None


def _get_layout_by_matching_name(prs: Presentation, targets: List[str]):
    target_set = {t.strip().lower() for t in targets}
    for layout in prs.slide_layouts:
        matching_name = (layout._element.get("matchingName") or "").strip().lower()
        if matching_name in target_set:
            return layout
    return None


def _get_layout_by_name(prs: Presentation, targets: List[str]):
    target_set = {t.strip().lower() for t in targets}
    for layout in prs.slide_layouts:
        if (layout.name or "").strip().lower() in target_set:
            return layout
    return None


def _get_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").strip()


def _find_layout_anchor(layout):
    anchors: Dict[str, Any] = {}
    for shp in layout.shapes:
        txt = _get_text(shp).upper()
        if "RRP" in txt:
            anchors["rrp_label"] = shp
        if "COLORWAY" in txt:
            anchors["color_label"] = shp
    return anchors


def _add_text_by_spec(slide, text: str, spec: Dict[str, Any], color_override: str | None = None):
    tb = slide.shapes.add_textbox(
        Mm(spec["left"]),
        Mm(spec["top"]),
        Mm(spec["width"]),
        Mm(spec["height"]),
    )
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0

    run = p.add_run()
    run.text = text
    f = run.font
    f.name = spec["font_name"]
    f.size = Pt(spec["font_size"])
    f.bold = spec.get("bold")

    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", spec["font_name"])

    rgb = _hex_to_rgbcolor(color_override or spec.get("color_hex"))
    if rgb is not None:
        try:
            f.color.rgb = rgb
        except Exception:
            pass


def _add_text_at(slide, text: str, left_mm: float, top_mm: float, width_mm: float, height_mm: float):
    tb = slide.shapes.add_textbox(Mm(left_mm), Mm(top_mm), Mm(width_mm), Mm(height_mm))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.0

    run = p.add_run()
    run.text = text
    run.font.name = "Averta Light"
    run.font.size = Pt(10)
    run.font.bold = False

    rgb = _hex_to_rgbcolor("#000000")
    if rgb is not None:
        run.font.color.rgb = rgb

    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", "Averta Light")


def _format_color_name(name: Any) -> str:
    if not name:
        return ""
    return str(name).strip().upper()


def _strip_vendor_watermark(prs: Presentation):
    markers = ("VORLAGENBAUER", "ERSTELLT DURCH")
    targets = list(prs.slide_masters)
    for m in prs.slide_masters:
        targets.extend(list(m.slide_layouts))

    for container in targets:
        for shp in container.shapes:
            if not getattr(shp, "has_text_frame", False):
                continue
            txt = (shp.text or "").upper()
            if any(marker in txt for marker in markers):
                shp.text_frame.clear()


def _ensure_slide_number_enabled(prs: Presentation):
    hf = prs._element.find(qn("p:hf"))
    if hf is None:
        hf = OxmlElement("p:hf")
        prs._element.append(hf)
    hf.set("sldNum", "1")


def _load_artwork_meta(meta_path: str) -> Dict[str, str]:
    if not os.path.exists(meta_path):
        return {}
    try:
        with open(meta_path, "r", encoding="utf-8") as f:
            data = json.load(f)
            return data if isinstance(data, dict) else {}
    except Exception:
        return {}


def _get_artwork_mode(art_name: str, meta: Dict[str, str]) -> str:
    mode = meta.get(art_name, ARTWORK_MODE_DEFAULT)
    if mode not in (ARTWORK_MODE_DEFAULT, ARTWORK_MODE_HORIZONTAL, ARTWORK_MODE_SMALL):
        return ARTWORK_MODE_DEFAULT
    return mode


def generate_pptx(
    products: List[Dict[str, Any]],
    template_file: str = "template.pptx",
    logo_dir: str = "assets/logos",
    artwork_dir: str = "assets/artworks",
):
    prs = Presentation(template_file) if os.path.exists(template_file) else Presentation()
    _strip_vendor_watermark(prs)
    _ensure_slide_number_enabled(prs)

    selected_layout = (
        _get_layout_by_matching_name(prs, ["default"]) 
        or _get_layout_by_matching_name(prs, ["title"]) 
        or _get_layout_by_name(prs, ["HB Title / Content", "CUSTOM"])
    )
    if selected_layout is None:
        selected_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    layout_anchors = _find_layout_anchor(selected_layout)
    artwork_meta = _load_artwork_meta(os.path.join(artwork_dir, "_meta.json"))

    for data in products:
        slide = prs.slides.add_slide(selected_layout)

        season_name = data.get("season_item", "")
        if season_name:
            _add_text_by_spec(slide, season_name, TEXT_SPECS["season"], color_override=data.get("season_color"))
        _add_text_by_spec(slide, data.get("name", ""), TEXT_SPECS["category"])
        _add_text_by_spec(slide, data.get("code", ""), TEXT_SPECS["code"])

        if data.get("rrp"):
            rrp_left = layout_anchors["rrp_label"].left if layout_anchors.get("rrp_label") else Mm(250)
            rrp_top = layout_anchors["rrp_label"].top if layout_anchors.get("rrp_label") else Mm(15)
            rrp = slide.shapes.add_textbox(rrp_left, rrp_top, Mm(50), Mm(15))
            rrp.text_frame.text = f"RRP : {data['rrp']}"
            rrp.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        if data.get("main_image"):
            main_pic = slide.shapes.add_picture(data["main_image"], left=Mm(0), top=Mm(0), width=Mm(MAIN_IMAGE_WIDTH_MM))
            main_pic.left = int(Mm(MAIN_IMAGE_CENTER_X_MM) - (main_pic.width / 2))
            main_pic.top = int(Mm(MAIN_IMAGE_CENTER_Y_MM) - (main_pic.height / 2))

        if data.get("logo") and data["logo"] != "선택 없음":
            p_logo = os.path.join(logo_dir, data["logo"])
            if os.path.exists(p_logo):
                logo_pic = slide.shapes.add_picture(p_logo, left=Mm(0), top=Mm(0), height=Mm(LOGO_HEIGHT_MM))
                logo_pic.left = int(Mm(LOGO_CENTER_X_MM) - (logo_pic.width / 2))
                logo_pic.top = int(Mm(LOGO_CENTER_Y_MM) - (logo_pic.height / 2))

        artworks = data.get("artworks", [])
        if artworks:
            current_top = int(Mm(ARTWORK_START_TOP_MM))
            gap_emu = int(Mm(ARTWORK_VERTICAL_GAP_MM))
            for art_name in artworks:
                p_art = os.path.join(artwork_dir, art_name)
                if not os.path.exists(p_art):
                    continue
                mode = _get_artwork_mode(art_name, artwork_meta)
                if mode == ARTWORK_MODE_SMALL:
                    art_pic = slide.shapes.add_picture(p_art, left=Mm(0), top=Mm(0), width=Mm(ARTWORK_SMALL_WIDTH_MM))
                elif mode == ARTWORK_MODE_HORIZONTAL:
                    art_pic = slide.shapes.add_picture(p_art, left=Mm(0), top=Mm(0), width=Mm(ARTWORK_DEFAULT_WIDTH_MM))
                else:
                    art_pic = slide.shapes.add_picture(p_art, left=Mm(0), top=Mm(0), height=Mm(ARTWORK_PORTRAIT_HEIGHT_MM))
                art_pic.left = int(Mm(ARTWORK_CENTER_X_MM) - (art_pic.width / 2))
                art_pic.top = current_top
                current_top += art_pic.height + gap_emu

        sx, sy, w, g = 180, COLORWAY_IMAGE_TOP_MM, COLORWAY_IMAGE_WIDTH_MM, 5
        colors = data.get("colors", [])
        color_count = len(colors)
        is_two = color_count == 2
        is_three = color_count == 3
        if is_two:
            sx = COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM
            sy = COLORWAY_IMAGE_TOP_MM
        if is_three:
            sx = COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM
            sy = COLORWAY_IMAGE_TOP_MM

        per_row = 3
        row_gap = 8
        img_h = 30
        rows = max(1, (len(colors) + per_row - 1) // per_row) if colors else 1
        circled_nums = ["①", "②", "③", "④", "⑤", "⑥", "⑦", "⑧", "⑨", "⑩"]
        for i, c in enumerate(colors):
            row = i // per_row
            col = i % per_row
            cy = sy - (rows - 1 - row) * (img_h + row_gap + 10)
            cx = sx + (col * (w + g))
            if is_two:
                cx = COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_TWO_ITEMS_LABEL_GAP_MM)
                cy = COLORWAY_IMAGE_TOP_MM
            if is_three:
                cx = COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM + (i * COLORWAY_THREE_ITEMS_LABEL_GAP_MM)
                cy = COLORWAY_IMAGE_TOP_MM

            if c.get("img"):
                slide.shapes.add_picture(c["img"], left=Mm(cx), top=Mm(cy), width=Mm(COLORWAY_IMAGE_WIDTH_MM))

            label = f"{circled_nums[i]}{_format_color_name(c.get('name'))}"
            if (is_two or is_three) and rows == 1:
                start_left = COLORWAY_TWO_ITEMS_LABEL_START_LEFT_MM if is_two else COLORWAY_THREE_ITEMS_LABEL_START_LEFT_MM
                start_top = COLORWAY_TWO_ITEMS_LABEL_TOP_MM if is_two else COLORWAY_THREE_ITEMS_LABEL_TOP_MM
                gap = COLORWAY_TWO_ITEMS_LABEL_GAP_MM if is_two else COLORWAY_THREE_ITEMS_LABEL_GAP_MM
                _add_text_at(slide, label, start_left + (i * gap), start_top, 32.0, 5.0)
            else:
                tb = slide.shapes.add_textbox(Mm(cx), Mm(cy + img_h + 2), Mm(w), Mm(10))
                tb.text_frame.text = _format_color_name(c.get("name"))
                tb.text_frame.paragraphs[0].font.size = Pt(9)
                tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
